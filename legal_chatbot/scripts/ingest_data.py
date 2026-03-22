import os
import json
import re

import fitz 
from docx import Document
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

DATA_DIR = "data"
OUTPUT_DIR = os.path.join("data", "processed")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def get_all_files(root_dir: str):
    files = []
    for current_dir, _, filenames in os.walk(root_dir):
        for name in filenames:
            full_path = os.path.join(current_dir, name)
            files.append(full_path)
    return files


# ---------- Normalization helpers ----------

def clean_ocr_text(text: str) -> str:
    """
    Advanced OCR cleaning for PDF extractions.
    Removes noise, fixes spacing, normalizes punctuation, ensures readability.
    """
    if not text:
        return ""

    # Remove Hindi/non-ASCII characters (keep only English text)
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)
    
    # Remove standalone page numbers and isolated digits on their own lines
    text = re.sub(r'\n\s*\d{1,4}\s*\n', '\n', text)
    text = re.sub(r'^\s*\d{1,4}\s*$', '', text, flags=re.MULTILINE)
    
    # Remove empty brackets and parentheses
    text = re.sub(r'\(\s*\)', '', text)
    text = re.sub(r'\[\s*\]', '', text)
    text = re.sub(r'\{\s*\}', '', text)
    
    # Remove broken patterns like "( ) , ," or "0 0 0 0"
    text = re.sub(r'\(\s*\)\s*[,;.]*\s*', '', text)
    text = re.sub(r'(\d\s+){3,}', ' ', text)  # Remove repeated spaced digits
    
    # Fix hyphenated words at line breaks
    text = re.sub(r'-\s*\n\s*', '', text)
    
    # Remove excessive punctuation (e.g., ",,," or "...")
    text = re.sub(r'[,;.]{3,}', '.', text)
    text = re.sub(r'[,;]{2,}', ',', text)
    
    # Remove lines with only punctuation/symbols
    text = re.sub(r'^\s*[^\w\s]+\s*$', '', text, flags=re.MULTILINE)
    
    # Normalize whitespace
    text = re.sub(r' {2,}', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'\t+', ' ', text)
    
    # Remove trailing/leading whitespace from each line
    lines = [line.strip() for line in text.split('\n')]
    
    # Remove empty lines and very short meaningless lines
    cleaned_lines = []
    for line in lines:
        if len(line) > 2 and re.search(r'\w', line):
            cleaned_lines.append(line)
    
    text = '\n'.join(cleaned_lines)
    
    # Final cleanup: ensure sentences are properly spaced
    text = re.sub(r'\.([A-Z])', r'. \1', text)
    text = re.sub(r',([^\s])', r', \1', text)
    
    return text.strip()


def normalize_text(text: str) -> str:
    """Apply comprehensive OCR cleaning."""
    return clean_ocr_text(text)


def remove_repeating_headers_footers(pages_text: list[str]) -> str:
    """
    Very simple heuristic:
    - split text per page
    - find lines that appear in many pages (e.g., case title, court name, page numbers)
    - remove those lines from all pages
    """
    line_counts = {}

    # collect line frequencies
    for page in pages_text:
        lines = set(page.split("\n"))  # set -> count each unique line once per page
        for line in lines:
            line = line.strip()
            if not line:
                continue
            line_counts[line] = line_counts.get(line, 0) + 1

    if not pages_text:
        return ""

    threshold = max(2, int(0.5 * len(pages_text)))  # appears in >= 50% of pages
    repeating_lines = {line for line, count in line_counts.items() if count >= threshold}

    cleaned_pages = []
    for page in pages_text:
        new_lines = []
        for line in page.split("\n"):
            if line.strip() in repeating_lines:
                continue
            new_lines.append(line)
        cleaned_pages.append("\n".join(new_lines))

    return "\n\n".join(cleaned_pages)


# ---------- Extraction functions ----------

def extract_pdf(path: str):
    print("Extracting PDF:", path)
    pages = []
    meta = {"page_count": 0, "title": None}

    try:
        doc = fitz.open(path)
        meta["page_count"] = doc.page_count
        meta["title"] = doc.metadata.get("title")
        for page in doc:
            pages.append(page.get_text("text"))
        doc.close()
    except Exception as e:
        print("Error reading PDF:", e)
        return "", meta

    # remove headers/footers before global normalization
    text_no_headers = remove_repeating_headers_footers(pages)
    normalized = normalize_text(text_no_headers)
    return normalized, meta


def extract_docx(path: str):
    print("Extracting DOCX:", path)
    text = ""
    meta = {"page_count": None, "title": None}

    try:
        doc = Document(path)
        # python-docx doesn't expose pages; treat as one big text
        for p in doc.paragraphs:
            text += p.text + "\n"

        # core properties -> title (may be empty)
        core = doc.core_properties
        meta["title"] = core.title or None
    except Exception as e:
        print("Error reading DOCX:", e)
        return "", meta

    normalized = normalize_text(text)
    return normalized, meta


def extract_pptx(path: str):
    print("Extracting PPTX:", path)
    text = ""
    meta = {"slide_count": 0, "title": None}

    try:
        prs = Presentation(path)
        meta["slide_count"] = len(prs.slides)
        meta["title"] = prs.core_properties.title or None

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print("Error reading PPTX:", e)
        return "", meta

    normalized = normalize_text(text)
    return normalized, meta


def extract_html(path: str):
    print("Extracting HTML:", path)
    text = ""
    meta = {"title": None}

    try:
        with open(path, "r", encoding="utf-8") as f:
            html = f.read()

        # crude title extraction
        m = re.search(r"<title>(.*?)</title>", html, flags=re.IGNORECASE | re.DOTALL)
        if m:
            meta["title"] = re.sub(r"\s+", " ", m.group(1)).strip()

        # strip tags
        text = re.sub(r"<[^>]+>", " ", html)
    except Exception as e:
        print("Error reading HTML:", e)
        return "", meta

    normalized = normalize_text(text)
    return normalized, meta


def chunk_text(raw_text: str):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=2000,
        chunk_overlap=200,
    )
    return splitter.split_text(raw_text)


def is_valid_legal_chunk(text: str) -> bool:
    """Filter out chunks that are not useful legal content."""
    text_lower = text.lower()
    
    # Minimum length requirement - must have substantial content
    if len(text.strip()) < 200:
        return False
    
    # Skip chunks that contain educational/navigation patterns
    skip_patterns = [
        "after studying this lesson",
        "you will be able to",
        "learning objectives",
        "intext questions",
        "fill in the blanks",
        "explore dla piper",
        "related resources",
        "click here",
        "table of contents",
        "contents page",
        "index page",
        "© copyright",
        "all rights reserved",
        "terms and conditions",
        "privacy policy",
        "cookie policy",
        "new rights",
        "more information",
        "read more",
        "learn more",
        "visit our",
        "subscribe",
        "sign up",
    ]
    
    for pattern in skip_patterns:
        if pattern in text_lower:
            return False
    
    # Skip if text starts with "Objectives" (educational material)
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if lines and lines[0].lower().startswith("objectives"):
        return False
    
    # Skip if chunk is just a title/heading (very short with no substance)
    if len(lines) <= 2 and len(text) < 300:
        return False
    
    # Must contain some legal-related terms
    legal_terms = [
        "section", "article", "act", "law", "court", "legal", "rights",
        "constitution", "punishment", "offence", "crime", "penalty",
        "shall", "person", "property", "contract", "liability", "india",
        "imprisonment", "fine", "guilty", "accused", "evidence", "trial",
        "plaintiff", "defendant", "judge", "justice", "statute", "citizen",
        "government", "parliament", "state", "provision", "clause", "amendment"
    ]
    
    # Count how many legal terms are present
    legal_term_count = sum(1 for term in legal_terms if term in text_lower)
    
    # Must have at least 3 legal terms for quality content
    return legal_term_count >= 3


def main():
    all_chunks = []
    manifest = []

    print("\nReading files from data folder (including subfolders)...\n")

    all_files = get_all_files(DATA_DIR)

    for path in all_files:
        # skip processed output folder
        if os.path.abspath(OUTPUT_DIR) in os.path.abspath(path):
            continue

        file = os.path.basename(path)
        print("Checking file:", path)

        ext = file.lower()

        text = ""
        meta = {}

        if ext.endswith(".pdf"):
            text, meta = extract_pdf(path)
            file_type = "pdf"
        elif ext.endswith(".docx"):
            text, meta = extract_docx(path)
            file_type = "docx"
        elif ext.endswith(".pptx"):
            text, meta = extract_pptx(path)
            file_type = "pptx"
        elif ext.endswith(".html") or ext.endswith(".htm"):
            text, meta = extract_html(path)
            file_type = "html"
        elif ext.endswith(".txt"):
            print("Extracting TXT:", path)
            try:
                with open(path, "r", encoding="utf-8") as f:
                    text = f.read()
                meta = {"title": None}
                text = normalize_text(text)
                file_type = "txt"
            except Exception as e:
                print("Error reading TXT:", e)
                continue
        else:
            print("Skipping unsupported file:", file)
            continue

        if not text.strip():
            print("No text extracted from:", file)
            continue

        print("Chunking text...")
        chunks = chunk_text(text)
        print("Chunks created:", len(chunks))

        # add to manifest (per-document metadata)
        manifest.append(
            {
                "source_file": file,
                "source_path": path,
                "file_type": file_type,
                "num_chunks": len(chunks),
                **meta,
            }
        )

        for i, c in enumerate(chunks):
            clean_text = c.strip()
            
            # Filter out invalid/irrelevant chunks
            if not is_valid_legal_chunk(clean_text):
                continue

            all_chunks.append(
                {
                    "source_file": file,
                    "source_path": path,
                    "file_type": file_type,
                    "chunk_id": i,
                    "chunk_length": len(clean_text),
                    "chunk_text": clean_text,
                }
            )

    if not all_chunks:
        print("\nNo chunks created. Check your files.")
        return

    # Save chunk JSON
    json_out = os.path.join(OUTPUT_DIR, "legal_chunks.json")
    with open(json_out, "w", encoding="utf-8") as f:
        json.dump(all_chunks, f, indent=4, ensure_ascii=False)

    # Save manifest
    manifest_out = os.path.join(OUTPUT_DIR, "manifest.json")
    with open(manifest_out, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=4, ensure_ascii=False)

    print("\nDone.")
    print("JSON chunks saved at:", json_out)
    print("Manifest saved at:", manifest_out)
    print("Total chunks:", len(all_chunks), "\n")


if __name__ == "__main__":
    main()
